"""Präsentation-Agent: Zwischen- und Schlusspräsentation als .pptx."""
from __future__ import annotations
import asyncio
import json
import time
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from scripts import config
from scripts.coherence import Universe
from scripts.subagents.base import (
    SubagentResult, claude_opus_complete, render_prompt,
    emit_start, emit_done,
)
from scripts.utils import log

PHASE = 4
NAME = "praesentation"
SYSTEM = "Du planst Präsentationen. Du antwortest nur mit reinem JSON."

ZAG_MAGENTA = RGBColor(0xE3, 0x00, 0x59)
ZAG_LILA = RGBColor(0x7A, 0x00, 0xDF)
BLACK = RGBColor(0x0A, 0x0A, 0x0F)


def _extract_json(raw: str) -> dict:
    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1].rsplit("```", 1)[0]
    return json.loads(raw)


def _build_pptx(slides: list[dict], title_slide_text: dict, out: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Titel-Slide
    blank_layout = prs.slide_layouts[6]
    s0 = prs.slides.add_slide(blank_layout)
    _add_background(s0, prs, BLACK)
    _add_text(s0, title_slide_text["titel"], Inches(0.8), Inches(2.5), Inches(12), Inches(2), Pt(44), True, ZAG_MAGENTA)
    _add_text(s0, title_slide_text.get("sub", ""), Inches(0.8), Inches(4.8), Inches(12), Inches(1), Pt(22), False, RGBColor(0xF0, 0xF0, 0xF0))

    # Content-Slides
    for sl in slides:
        s = prs.slides.add_slide(blank_layout)
        _add_background(s, prs, RGBColor(0xF5, 0xF5, 0xF5))
        _add_text(s, sl.get("titel", ""), Inches(0.8), Inches(0.5), Inches(12), Inches(1), Pt(32), True, ZAG_LILA)
        bullets = sl.get("bullets", [])
        tx = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(5)).text_frame
        tx.word_wrap = True
        for i, b in enumerate(bullets):
            p = tx.paragraphs[0] if i == 0 else tx.add_paragraph()
            p.text = "• " + b
            for run in p.runs:
                run.font.size = Pt(22)
                run.font.color.rgb = BLACK
        # Sprechnotizen
        notes = s.notes_slide.notes_text_frame
        notes.text = sl.get("sprechnotizen", "")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def _add_background(slide, prs, color: RGBColor):
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def _add_text(slide, text, left, top, width, height, font_size, bold, color):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


async def run(u: Universe) -> SubagentResult:
    await emit_start(NAME, PHASE)
    t0 = time.monotonic()

    zw_task = claude_opus_complete(system=SYSTEM,
        user=render_prompt("praesentation_slides.j2", u=u, variante="zwischen"),
        max_tokens=3500)
    sch_task = claude_opus_complete(system=SYSTEM,
        user=render_prompt("praesentation_slides.j2", u=u, variante="schluss"),
        max_tokens=6000)
    zw_json, sch_json = await asyncio.gather(zw_task, sch_task)

    zw = _extract_json(zw_json)
    sch = _extract_json(sch_json)

    zw_pptx = config.ARTIFACTS_DIR / "18_zwischenpraesentation.pptx"
    sch_pptx = config.ARTIFACTS_DIR / "19_schlusspraesentation.pptx"
    _build_pptx(
        zw["slides"],
        {"titel": f"Zwischenstand: {u.thema.titel}", "sub": f"{u.schuelerin.vorname} {u.schuelerin.nachname} · {u.schuelerin.klasse}"},
        zw_pptx,
    )
    _build_pptx(
        sch["slides"],
        {"titel": u.thema.titel, "sub": f"Vertiefungsarbeit · {u.schuelerin.vorname} {u.schuelerin.nachname}"},
        sch_pptx,
    )

    # Sprechnotizen als Markdown extrahieren
    notes_md = f"# Sprechnotizen Schlusspräsentation · {u.schuelerin.vorname} {u.schuelerin.nachname}\n\n"
    for i, sl in enumerate(sch["slides"], 1):
        notes_md += f"## Slide {i}: {sl.get('titel', '')}\n\n{sl.get('sprechnotizen', '')}\n\n"
    (config.ARTIFACTS_DIR / "19_schlusspraesentation_sprechnotizen.md").write_text(notes_md, encoding="utf-8")

    duration = time.monotonic() - t0
    await emit_done(NAME, PHASE, detail=f"zwischen ({len(zw['slides'])}), schluss ({len(sch['slides'])})")
    log.info(f"[{NAME}] Präsentationen fertig in {duration:.1f}s")
    return SubagentResult(name=NAME, output_path=sch_pptx, duration_s=duration)
